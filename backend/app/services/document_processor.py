"""
Unified document processor supporting multiple file formats
"""
import os
import csv
import chardet
from typing import List, Dict, Optional, Tuple
from pathlib import Path
import structlog

# PDF
import fitz  # PyMuPDF

# DOCX
from docx import Document as DocxDocument

# Excel
from openpyxl import load_workbook

# PowerPoint
from pptx import Presentation

# Markdown
import markdown

logger = structlog.get_logger()


class DocumentProcessor:
    """Unified service for processing multiple document formats"""

    SUPPORTED_FORMATS = {
        'pdf': ['.pdf'],
        'docx': ['.docx', '.doc'],
        'excel': ['.xlsx', '.xls'],
        'csv': ['.csv'],
        'text': ['.txt', '.md', '.markdown', '.rst'],
        'powerpoint': ['.pptx', '.ppt'],
    }

    @classmethod
    def get_file_type(cls, file_path: str) -> str:
        """
        Detect file type from extension

        Args:
            file_path: Path to file

        Returns:
            File type ('pdf', 'docx', 'excel', 'csv', 'text', 'powerpoint', 'unknown')
        """
        ext = Path(file_path).suffix.lower()

        for file_type, extensions in cls.SUPPORTED_FORMATS.items():
            if ext in extensions:
                return file_type

        return 'unknown'

    @classmethod
    def extract_metadata(cls, file_path: str) -> Dict[str, Optional[str]]:
        """
        Extract metadata (author, title, etc.) from document

        Args:
            file_path: Path to document file

        Returns:
            dict: Metadata dict with keys like 'author', 'title', etc.
        """
        file_type = cls.get_file_type(file_path)
        metadata = {"author": None, "title": None}

        try:
            if file_type == 'pdf':
                doc = fitz.open(file_path)
                pdf_metadata = doc.metadata
                metadata["author"] = pdf_metadata.get("author") or pdf_metadata.get("creator")
                metadata["title"] = pdf_metadata.get("title")
                doc.close()

            elif file_type == 'docx':
                doc = DocxDocument(file_path)
                core_props = doc.core_properties
                metadata["author"] = core_props.author
                metadata["title"] = core_props.title

        except Exception as e:
            logger.warning("metadata_extraction_failed", file_path=file_path, error=str(e))

        return metadata

    @classmethod
    def extract_text(cls, file_path: str) -> Tuple[str, int, List[Dict]]:
        """
        Extract text from any supported document format

        Args:
            file_path: Path to document file

        Returns:
            tuple: (full_text, page_count, pages_data)
                - full_text: All text concatenated
                - page_count: Number of pages/sheets/sections
                - pages_data: List of dicts with page/section-level data
        """
        file_type = cls.get_file_type(file_path)

        logger.info(
            "document_extraction_started",
            file_path=file_path,
            file_type=file_type,
        )

        try:
            if file_type == 'pdf':
                return cls._extract_pdf(file_path)
            elif file_type == 'docx':
                return cls._extract_docx(file_path)
            elif file_type == 'excel':
                return cls._extract_excel(file_path)
            elif file_type == 'csv':
                return cls._extract_csv(file_path)
            elif file_type == 'text':
                return cls._extract_text_file(file_path)
            elif file_type == 'powerpoint':
                return cls._extract_pptx(file_path)
            else:
                raise ValueError(f"Unsupported file format: {Path(file_path).suffix}")

        except Exception as e:
            logger.error(
                "document_extraction_failed",
                file_path=file_path,
                file_type=file_type,
                error=str(e),
                exc_info=True,
            )
            raise

    @staticmethod
    def _extract_pdf(file_path: str) -> Tuple[str, int, List[Dict]]:
        """Extract text from PDF"""
        doc = fitz.open(file_path)
        page_count = len(doc)
        pages_data = []
        full_text = []

        for page_num, page in enumerate(doc, start=1):
            text = page.get_text("text")
            pages_data.append({
                "page_number": page_num,
                "text": text,
                "char_count": len(text),
            })
            full_text.append(text)

        doc.close()

        logger.info(
            "pdf_extracted",
            file_path=file_path,
            page_count=page_count,
        )

        return "\n\n".join(full_text), page_count, pages_data

    @staticmethod
    def _extract_docx(file_path: str) -> Tuple[str, int, List[Dict]]:
        """Extract text from DOCX"""
        doc = DocxDocument(file_path)

        # Extract all paragraphs
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)

        # Extract tables
        table_texts = []
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text for cell in row.cells])
                if row_text.strip():
                    table_texts.append(row_text)

        # Combine all text
        all_text = paragraphs + table_texts
        full_text = "\n\n".join(all_text)

        # Create pages_data (DOCX doesn't have pages, use sections)
        pages_data = [{
            "page_number": 1,
            "text": full_text,
            "char_count": len(full_text),
            "paragraph_count": len(paragraphs),
            "table_count": len(doc.tables),
        }]

        logger.info(
            "docx_extracted",
            file_path=file_path,
            paragraphs=len(paragraphs),
            tables=len(doc.tables),
        )

        return full_text, 1, pages_data

    @staticmethod
    def _extract_excel(file_path: str) -> Tuple[str, int, List[Dict]]:
        """Extract text from Excel"""
        wb = load_workbook(file_path, data_only=True)
        sheet_count = len(wb.worksheets)
        pages_data = []
        full_text = []

        for sheet_num, sheet in enumerate(wb.worksheets, start=1):
            sheet_text = []

            # Extract cell values
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip(" |"):
                    sheet_text.append(row_text)

            text = "\n".join(sheet_text)
            pages_data.append({
                "page_number": sheet_num,
                "sheet_name": sheet.title,
                "text": text,
                "char_count": len(text),
                "row_count": sheet.max_row,
                "col_count": sheet.max_column,
            })
            full_text.append(f"## Sheet: {sheet.title}\n{text}")

        wb.close()

        logger.info(
            "excel_extracted",
            file_path=file_path,
            sheet_count=sheet_count,
        )

        return "\n\n".join(full_text), sheet_count, pages_data

    @staticmethod
    def _extract_csv(file_path: str) -> Tuple[str, int, List[Dict]]:
        """Extract text from CSV"""
        # Detect encoding
        with open(file_path, 'rb') as f:
            raw_data = f.read()
            result = chardet.detect(raw_data)
            encoding = result['encoding'] or 'utf-8'

        # Read CSV
        rows = []
        with open(file_path, 'r', encoding=encoding) as f:
            reader = csv.reader(f)
            for row in reader:
                row_text = " | ".join(row)
                if row_text.strip(" |"):
                    rows.append(row_text)

        full_text = "\n".join(rows)

        pages_data = [{
            "page_number": 1,
            "text": full_text,
            "char_count": len(full_text),
            "row_count": len(rows),
        }]

        logger.info(
            "csv_extracted",
            file_path=file_path,
            row_count=len(rows),
        )

        return full_text, 1, pages_data

    @staticmethod
    def _extract_text_file(file_path: str) -> Tuple[str, int, List[Dict]]:
        """Extract text from plain text/markdown files"""
        # Detect encoding
        with open(file_path, 'rb') as f:
            raw_data = f.read()
            result = chardet.detect(raw_data)
            encoding = result['encoding'] or 'utf-8'

        # Read text
        with open(file_path, 'r', encoding=encoding) as f:
            full_text = f.read()

        # Check if it's markdown
        ext = Path(file_path).suffix.lower()
        is_markdown = ext in ['.md', '.markdown']

        # For markdown, we could parse it to HTML then extract text
        # But for search purposes, raw markdown is often better

        pages_data = [{
            "page_number": 1,
            "text": full_text,
            "char_count": len(full_text),
            "is_markdown": is_markdown,
        }]

        logger.info(
            "text_file_extracted",
            file_path=file_path,
            is_markdown=is_markdown,
            char_count=len(full_text),
        )

        return full_text, 1, pages_data

    @staticmethod
    def _extract_pptx(file_path: str) -> Tuple[str, int, List[Dict]]:
        """Extract text from PowerPoint"""
        prs = Presentation(file_path)
        slide_count = len(prs.slides)
        pages_data = []
        full_text = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = []

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)

            text = "\n".join(slide_text)
            pages_data.append({
                "page_number": slide_num,
                "text": text,
                "char_count": len(text),
                "shape_count": len(slide.shapes),
            })
            full_text.append(f"## Slide {slide_num}\n{text}")

        logger.info(
            "pptx_extracted",
            file_path=file_path,
            slide_count=slide_count,
        )

        return "\n\n".join(full_text), slide_count, pages_data

    @classmethod
    def get_metadata(cls, file_path: str) -> Dict:
        """
        Get document metadata

        Args:
            file_path: Path to file

        Returns:
            dict: Document metadata
        """
        file_type = cls.get_file_type(file_path)
        file_stat = os.stat(file_path)

        metadata = {
            "file_type": file_type,
            "file_name": Path(file_path).name,
            "file_size": file_stat.st_size,
            "file_extension": Path(file_path).suffix,
        }

        try:
            if file_type == 'pdf':
                doc = fitz.open(file_path)
                pdf_meta = doc.metadata
                metadata.update({
                    "title": pdf_meta.get("title", ""),
                    "author": pdf_meta.get("author", ""),
                    "subject": pdf_meta.get("subject", ""),
                    "page_count": len(doc),
                })
                doc.close()
            elif file_type == 'docx':
                doc = DocxDocument(file_path)
                core_props = doc.core_properties
                metadata.update({
                    "title": core_props.title or "",
                    "author": core_props.author or "",
                    "subject": core_props.subject or "",
                })
            elif file_type == 'excel':
                wb = load_workbook(file_path, data_only=True)
                metadata["sheet_count"] = len(wb.worksheets)
                metadata["sheet_names"] = [s.title for s in wb.worksheets]
                wb.close()

        except Exception as e:
            logger.warning(
                "metadata_extraction_failed",
                file_path=file_path,
                error=str(e),
            )

        return metadata
